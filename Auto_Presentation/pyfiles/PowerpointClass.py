from pptx import Presentation
from pptx.util import Inches
import imaginairy
from imaginairy.schema import ImaginePrompt
import random

class PowerPointCreator:
    def __init__(self):
        self.presentation = Presentation()
        self.filename=""

    def add_title_slide(self, title, subtitle):
        # Add a title slide
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.placeholders[0]
        subtitle_placeholder = slide.placeholders[1]
        title_placeholder.text = title
        subtitle_placeholder.text = subtitle

    def AI_generate(self,prom):
        prompt = ImaginePrompt(
                prompt=prom, 
                seed=random.randint(0, 1000), 
                model_weights="sd15", 
                size="hd"
            )
        self.filename=imaginairy.imagine_image_files(prompts=prompt,outdir="images")[0]
    def add_content_slide(self, title, content, image_path=None):
        # Add a content slide with a title, content, and an optional image
        slide_layout = self.presentation.slide_layouts[1] # Assuming layout 1 has a placeholder for an image
        slide = self.presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.placeholders[0] # Placeholder for the title
        content_placeholder = slide.placeholders[1] # Placeholder for the content

        # Set the title to the most left part
        title_placeholder.text = title
        title_text_frame = title_placeholder.text_frame
        title_text_frame.paragraphs[0].alignment = 1 # Align the title to the left

        # Set the content (assuming it's in a text box)
        content_placeholder.text = content
        content_text_frame = content_placeholder.text_frame
        content_text_frame.paragraphs[0].alignment = 1 # Align the content to the left

        # Add the image to the bottom right part
        if image_path:
            # Calculate the position for the bottom right corner
            image_width = Inches(1.81102)
            image_height = Inches(1.81102)
            image_left = Inches(8.188976) # Adjust the left position to ensure it's in the corner
            image_top = Inches(5.6889764)  # Adjust the top position to ensure it's in the corner

            # Add a picture shape to the slide with the calculated size and position
            image_shape = slide.shapes.add_picture(image_path, image_left, image_top, width=image_width, height=image_height)

        # Set the font to "Times New Roman" for the content slide
        title_text_frame.paragraphs[0].font.name = "Times New Roman"
        content_text_frame.paragraphs[0].font.name = "Times New Roman"






    def save_presentation(self, path):
        # Save the presentation
        self.presentation.save(path)

    def create(self, title, subtitle, image_path, content_list, path,AI_pic=True):
        # Create the presentation
        self.add_title_slide(title, subtitle)
        for i, content in enumerate(content_list, start=2):
            if AI_pic:
                self.AI_generate(content['content'])
                self.add_content_slide(f"{i}. {content['title']}", content['content']+'\naccording to '+content['references'],image_path=self.filename)
            else:
                self.add_content_slide(f"{i}. {content['title']}", content['content']+'\naccording to '+content['references'])
        self.save_presentation(path)
        print(f"Presentation created at {path}")


